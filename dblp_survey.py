import json as js
import os
import webbrowser
from os.path import join

import pandas as pd
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt
from requests_html import HTMLSession
from tqdm import tqdm
import math
from get import get_bibtex, get_ids, get_title, get_url
from plugin import (cross_modal_except_title, fine_tun_except_title,
                    multi_modal_except_title, pre_train_except_title,
                    summar_except_title)

headers = {
    'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_14_2) AppleWebKit/537.36 (KHTML, like Gecko) '
    'Chrome/71.0.3578.98 Safari/537.36'
}
url_prefix = ['https://dblp.org/search?q=', '%20year%3A', '%3A']
url_next_prefix = [
    'https://dblp.org/search/publ/inc?q=',
    '%20year%3A',
    '%3A',
    '&s=ydvspc&h=30&b=',
]

base_xpath = ['//*[@id="', '"]']
session = HTMLSession()
text_length = 90


def content_parse(url, headers, survey, exist_titles, res=None):
    while True:
        try:
            r = session.get(url, headers=headers)
        except:
            continue
        break
    # # Get Paper IDs
    raw_ids, ids = get_ids(r, survey)

    # # Get Paper Title / URL / Bibtex
    with tqdm(total=len(ids)) as pbar:
        for id in ids:
            base = f'{base_xpath[0]}{id}{base_xpath[1]}'
            # # # Get Paper Title
            title = get_title(r, base)
            if title.lower() in exist_titles or (res(title) if res else False):
                # print(title)
                del survey[id]
                pbar.update(1)
                continue
            # # # Get Paper URL
            paper_url = get_url(r, base)
            # # # Get Paper Bibtex
            bibtex = get_bibtex(id, session, headers)

            survey[id]['title'] = title
            exist_titles.append(title.lower())
            survey[id]['paper_url'] = paper_url
            survey[id]['bibtex'] = bibtex
            pbar.update(1)
    r.close()
    return raw_ids, exist_titles


def make_pptx(query, year, survey):
    # Create Presentation
    left =  Cm(1)    # left，top为相对位置
    top =  Cm(0)
    top_url = Cm(2)
    top_bibtex = Cm(3)
    width = Cm(20)    # width，height为文本框的大小
    height = Cm(2)   
    url_height = Cm(1)
    bibtex_height = Cm(9)
    prs = Presentation()
    for id in survey:
        # Use the fourth template Create Slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # for shape in slide.placeholders:
        #     phf = shape.placeholder_format
        #     print(f'{phf.idx}--{shape.name}--{phf.type}')
        
        title = slide.shapes.add_textbox(left, top, width, height)
        tf = title.text_frame
        para = tf.add_paragraph()    # 新增段落
        title = survey[id]['title'].split(' ')
        titles = []
        length = 0
        for item in title:
            length += len(item)
            if length < text_length:
                titles.append(item)
            else:
                length = len(item)
                titles.append('\n')
                titles.append(item)
        title = ' '.join(titles)
        
        
        
        para.text = title+' ('+survey[id]['type'] + ')'  # 向段落写入文字
        para.alignment = PP_ALIGN.LEFT    # 居中
        para.line_spacing = 0.8    # 1.5 倍的行距c

        font = para.font
        font.name = 'Calibri'    # 字体类型
        font.bold = True    # 加粗
        font.size = Pt(15)    # 大小
        
        if survey[id]['paper_url']:
            try:
                url = slide.shapes.add_textbox(left, top_url, width, url_height)
                url_tf = url.text_frame
                url_para = url_tf.add_paragraph()    # 新增段落
                
                run = url_para.add_run()
                run.text = survey[id]['paper_url']
                run.hyperlink.address = survey[id]['paper_url']
                
                font = url_para.font
                font.name = 'Calibri'    # 字体类型
                font.bold = False    # 加粗
                font.size = Pt(15)    # 大小
                
            except Exception as e:
                print(repr(e))
        
        try:
            bibtex = slide.shapes.add_textbox(left, top_bibtex, width, bibtex_height)
            bibtex_tf = bibtex.text_frame
            bibtex_para = bibtex_tf.add_paragraph()    # 新增段落
            bibtex_para.line_spacing = 0.8    # 1.5 倍的行距c
            bibtex_para.text = survey[id]['bibtex']
            
            font = bibtex_para.font
            font.name = 'Calibri'    # 字体类型
            font.bold = False    # 加粗
            font.size = Pt(12)    # 大小
            
        except Exception as e:
            print(repr(e))
        
    prs.save(f'.\\save\\ppt\\{query}_{year}_dblp.pptx')


def request_url(query, year=None, res=None):
    if ' ' in query:
        query = query.replace(' ', '+')
    print(f'{query}_{year}')

    print(f'Start Get Web Page Content')
    survey = {}
    exist_titles = []
    # Get Main Web Page Content
    print('Get Main Web Page Content')
    if year:
        url = f'{url_prefix[0]}{query}{url_prefix[1]}{year}{url_prefix[2]}'
    else:
        url = f'{url_prefix[0]}{query}'
        
    _, exist_titles = content_parse(url, headers, survey, exist_titles, res)

    # Get Dynamic Web Page Content
    b = 1
    while True:
        print(f'Get Dynamic Web Page {b} Content')
        if year:
            next_url = f'{url_next_prefix[0]}{query}{url_next_prefix[1]}{year}{url_next_prefix[2]}{url_next_prefix[3]}{b}'
        else:
            next_url = f'{url_next_prefix[0]}{query}{url_next_prefix[3]}{b}'
            
        headers['Referer'] = url
        raw_ids, exist_titles = content_parse(next_url, headers, survey, exist_titles, res)
        b += 1
        if len(raw_ids) - 3 < 30:
            break
        # break

    print(f'Web Page Content Done: ', len(survey))

    print(f'Save Survey as Json')
    with open(f'.\\save\\json\\{query}_{year}_dblp.json', 'w') as file:
        js.dump(survey, file)
    # with open(f'./paper/save/{query}_{year}.json', 'r') as file:
    #     survey = js.load(file)

    print(f'Start Make PPT')
    make_pptx(query, year, survey)


def convert_ppt2excel(path):
    filelist = os.listdir(path)
    filelist = list(filter(lambda file: file.endswith('.pptx'), filelist))

    for filename in filelist:
        ptx = Presentation(join(path, filename))

        # df = pd.read_excel(join(path, filename.replace('pptx', 'xlsx')))
        data = {'title': [], 'type': [], 'paper_url': []}
        for slide in ptx.slides:
            type = slide.placeholders[1].text
            # if type == 'arXiv':
            #     break
            data['type'].append(slide.placeholders[1].text)
            data['title'].append(slide.placeholders[0].text)
            data['paper_url'].append(slide.placeholders[3].text)
        pd.DataFrame(data).to_excel(
            join(path, filename.replace('.pptx', '.xlsx')),
            columns=['title', 'type', 'paper_url'],
            index=False,
        )


def open_paper_page(path, arxiv=False):
    ptx = Presentation(path)

    for slide in ptx.slides:
        type = slide.placeholders[1].text
        if not arxiv and type == 'arXiv':
            continue
        paper_url = slide.placeholders[3].text
        webbrowser.open(paper_url)


def convert_ppt2latex(path):
    ptx = Presentation(path)

    bibtex_file = open(path.replace('.pptx', '_bibtex.txt'), 'w')
    summary_file = open(path.replace('.pptx', '_summary.txt'), 'w')
    
    for slide in ptx.slides:
        # type = slide.placeholders[1].text
        # if type == 'arXiv':
        #     break
        title = slide.placeholders[0].text
        # paper_url = slide.placeholders[3].text
        summary = slide.placeholders[2].text.replace('&', '\\&').replace('%', '\\%')
        bibtex = slide.placeholders[4].text
        bibtex_file.write(bibtex)
        # print(bibtex.split('\n')[1].split('{')[1].split(',')[0])
        cite = bibtex.split('\n')[1].split('{')[1].split(',')[0]
        try:
            summary_file.write(
                f'\n% {title}\n\\noindent \\textbf{{\\citet{{{cite}}}}}:\n{summary}\n'
            )
        except:
            print(cite)
            summary_file.write(f'\n\\noindent \\textbf{{\\citet{{{cite}}}}}:\n')
        # break


if __name__ == '__main__':
    # query = 'summar'
    # year = '2016'
    # request_url(query, year, summar_except_title)

    # query = 'multimodal document'
    # request_url(query)

    # query = 'multi-modal document'
    # request_url(query, res=multi_modal_except_title)

    # query = 'document understanding'
    # request_url(query)

    # query = 'document representation'
    # request_url(query)

    # query = 'visual rich document'
    # request_url(query)

    # query = 'pretrain'
    # year = '2019'
    # request_url(query, year)

    # query = 'pre-train'
    # year = '2019'
    # request_url(query, year, pre_train_except_title)

    # query = 'finetun'
    # year = '2019'
    # request_url(query, year)

    # query = 'fine-tun'
    # year = '2019'
    # request_url(query, year, fine_tun_except_title)

    # query = 'contrastive'
    # year = '2021'
    # request_url(query, year)

    # query = 'transformer'
    # year = '2021'
    # request_url(query, year)

    # query = 'neural architecture search'
    # year = '2021'
    # request_url(query, year)

    # query = 'relation caption'
    # request_url(query)

    # query = 'vision language'
    # years = ['2021', '2020', '2019']
    # for year in years:
    #     request_url(query, year)

    # query = 'multimodal'
    # years = ['2021', '2020', '2019']
    # for year in years:
    #     request_url(query, year)

    query = 'question answering'
    years = ['2021']
    for year in years:
        # request_url(query, year, multi_modal_except_title)
        request_url(query, year)

    # query = 'crossmodal'
    # years = ['2021', '2020', '2019']
    # for year in years:
    #     request_url(query, year)

    # query = 'cross-modal'
    # years = ['2021', '2020', '2019']
    # for year in years:
    #     request_url(query, year, cross_modal_except_title)

    # path = 'C:\\Users\\Haha\\Desktop\\survey\\summar\\summar_2016_add.pptx'
    # path = 'C:\\Users\\Haha\\Desktop\\survey\\multimodal document\\visual+rich+document_None.pptx'
    # path = 'C:/Users/Haha/Desktop/survey/image caption/relation+caption_None.pptx'
    # open_paper_page(path, arxiv=True)

    # path = 'C:\\Users\\Haha\\Desktop\\survey\\summar'
    # convert_ppt2excel(path)

    # path = 'C:\\Users\\Haha\\Desktop\\survey\\summar\\summar_2021_add.pptx'
    # path = 'C:\\Users\\Haha\\Desktop\\survey\\multimodal document\\visual+rich+document_None.pptx'
    # convert_ppt2latex(path)
