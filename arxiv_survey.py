import json as js
import multiprocessing as mp
import os
import pickle as pkl
import re
import time
import traceback
from os.path import join

import bs4
from matplotlib.pyplot import text
import requests
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt
from bs4 import BeautifulSoup
from lxml import etree
from rich import print
from selenium import webdriver  # 导入库
from selenium.webdriver.chrome.options import Options
from tqdm import tqdm


template_v1 = r'Showing (.*) of '
template_v2 = r'of (.*) results'
template_v3 = r'arXiv:(.*)'
max_papers = 100
text_length=90
summary_length = 100

def obtain_id_url(p):
    ''' 获取id，链接，pdf等信息'''
    arxiv_info = {}

    if p['class'][0] == 'list-title' and p['class'][1] == 'is-inline-block':
        for a in p.find_all(name='a'):
            if a.text == 'pdf':
                arxiv_info['pdf'] = a.get('href')
            elif a.text == 'other':
                arxiv_info['other'] = a.get('href')
            elif a.text[:5] == 'arXiv':
                arxiv_info['page'] = a.get('href')
                arxiv_info['id'] = re.findall(template_v3, a.text)[0]

    return arxiv_info


def obtain_author(p):
    ''' 获取作者信息'''
    author_info = []

    if p['class'][0] == 'authors':
        for a in p.find_all(name='a'):
            author_info.append(a.text)

    return author_info


def obtain_summary_info(p):
    ''' 获取摘要信息'''
    summary = ''

    if p['class'][0] == 'abstract' and p['class'][1] == 'mathjax':
        for span in p.find_all(name='span'):
            if span.attrs is not None and 'class' in span.attrs.keys() and span['class'][0] == 'abstract-full' and span['class'][1] == 'has-text-grey-dark' and span['class'][2] == 'mathjax':
                summary = span.text.split('\n')[1]

    return summary


def obtain_title(p):
    title = ''

    if p['class'][0] == 'title' and p['class'][1] == 'is-5' and p['class'][2] == 'mathjax':
        title = p.text.strip()

    return title


def obtain_comments(p):
    comment = ''

    if p['class'][0] == 'comments' and p['class'][1] == 'is-size-7':
        comment = p.text.replace('\n', ' ').replace('\t', ' ').replace('\r', ' ')

    return comment


def url_generation(query, search_type='all', order='-announced_date_first', size=50, start=0):
    query = query.lower().replace(' ', '+')
    url = f'https://arxiv.org/search/?query={query}&searchtype={search_type}&abstracts=show&order={order}&size={size}&start={start}'
    return url


def make_pptx(query, survey):
    # result.append({
    #                 'arxiv': arxiv_info,
    #                 'title': title,
    #                 'author_info': author_info,
    #                 'summary': summary,
    #                 'comment': comment
    #             })
    
    # Create Presentation
    left =  Cm(1)    # left，top为相对位置
    title_left = Cm(0)
    top =  Cm(0)
    top_url = Cm(2)
    top_summary = Cm(4)
    top_comment = Cm(9)
    width = Cm(20)    # width，height为文本框的大小
    height = Cm(2)   
    url_height = Cm(0.8)
    url_width = Cm(6.3)
    summary_height = Cm(6)
    
    prs = Presentation()
    for id in range(len(survey)):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(title_left, top, width, height)
        title_tf = title_box.text_frame
        para = title_tf.add_paragraph()    # 新增段落
        title = survey[id]['title'].split(' ')
        titles = []
        length = 0
        for item in title:
            length += len(item)
            if length < text_length:
                titles.append(item)
            else:
                length = len(item)
                titles.append('\n')
                titles.append(item)
        
        arxiv_id = survey[id]['arxiv']['id']
        titles.append(f'({arxiv_id})')
        
        title = ' '.join(titles)
        para.text = title  # 向段落写入文字
        para.alignment = PP_ALIGN.LEFT    # 居中
        para.line_spacing = 0.8    # 1.5 倍的行距c
        font = para.font
        font.name = 'Calibri'    # 字体类型
        font.bold = True    # 加粗
        font.size = Pt(15)    # 大小
        
        author = title_tf.add_paragraph()
        author.text = ','.join(survey[id]['author_info'])
        author.alignment = PP_ALIGN.LEFT    # 居中
        author.line_spacing = 0.8    # 1.5 倍的行距c
        font = author.font
        font.name = 'Calibri'    # 字体类型
        font.bold = False    # 加粗
        font.size = Pt(10)    # 大小
        
        try:
            url = slide.shapes.add_textbox(left, top_url, url_width, url_height)
            url_tf = url.text_frame
            url_para = url_tf.add_paragraph()    # 新增段落
            run = url_para.add_run()
            run.text = survey[id]['arxiv']['pdf'].strip()
            run.hyperlink.address = survey[id]['arxiv']['pdf']
            font = url_para.font
            font.name = 'Calibri'    # 字体类型
            font.bold = False    # 加粗
            font.size = Pt(10)    # 大小
        except Exception as e:
            print(repr(e))
    
        try:
            summary = slide.shapes.add_textbox(left, top_summary, width, summary_height)
            summary_tf = summary.text_frame
            summary_para = summary_tf.add_paragraph()    # 新增段落
            summary_para.line_spacing = 1.0    # 1.5 倍的行距c
            summary = survey[id]['summary'].split(' ')
            summarys = []
            length = 0
            for item in summary:
                length += len(item)
                if length < summary_length:
                    summarys.append(item)
                else:
                    length = len(item)
                    summarys.append('\n')
                    summarys.append(item)
                    
            summary = 'Abstract: '+' '.join(summarys)+'\n'
            summary_para.text = summary
            font = summary_para.font
            font.name = 'Calibri'    # 字体类型
            font.bold = False    # 加粗
            font.size = Pt(12)    # 大小
            
            comment_para = summary_tf.add_paragraph()    # 新增段落
            comment_para.text = survey[id]['comment']
            font = comment_para.font
            font.name = 'Calibri'    # 字体类型
            font.bold = False    # 加粗
            font.size = Pt(10)    # 大小
            
        except Exception as e:
            print(repr(e))
        
    prs.save(f'F:\\Survey\\save\\ppt\\{query}_arxiv.pptx')
    

def arxiv_survey(query, search_type='all', order='-announced_date_first', size=50, start=0):
    sum_num = 0
    chrome_options = Options()
    chrome_options.add_argument('--headless')
    driver = webdriver.Chrome(options=chrome_options)  # 声明浏览器
    result = []
    
    while True:
        input_url = url_generation(query, search_type=search_type, order=order, size=size, start=start)
        driver.get(input_url)
        time.sleep(7)
        soup = BeautifulSoup(driver.page_source, 'lxml')
        
        if sum_num == 0:
            for h1 in soup.find_all(name='h1'):
                if h1['class'][0] == 'title' and h1['class'][1] == 'is-clearfix':
                    h1_str = str(h1)
                    sum_num = int(re.findall(template_v2, h1_str)[0].replace(',', ''))

        for li in soup.find_all(name='li'):
            if li.attrs is not None and 'class' in li.attrs.keys() and li.attrs['class'][0] == 'arxiv-result':
                arxiv_info = {}
                title = ''
                author_info = []
                summary = ''
                comment = ''
                
                for p in li.find_all(name='p'):
                    if p['class'][0] == 'title' and p['class'][1] == 'is-5' and p['class'][2] == 'mathjax':
                        title = p.text.strip()
                    if p['class'][0] == 'list-title' and p['class'][1] == 'is-inline-block':
                        for a in p.find_all(name='a'):
                            if a.text == 'pdf':
                                arxiv_info['pdf'] = a.get('href')
                            elif a.text == 'other':
                                arxiv_info['other'] = a.get('href')
                            elif a.text[:5] == 'arXiv':
                                arxiv_info['page'] = a.get('href')
                                arxiv_info['id'] = re.findall(template_v3, a.text)[0]
                    if p['class'][0] == 'abstract' and p['class'][1] == 'mathjax':
                        for span in p.find_all(name='span'):
                            if span.attrs is not None and 'class' in span.attrs.keys() and span['class'][0] == 'abstract-full' and span['class'][1] == 'has-text-grey-dark' and span['class'][2] == 'mathjax':
                                summary = span.text.split('\n')[1].strip()
                    if p['class'][0] == 'authors':
                        for a in p.find_all(name='a'):
                            author_info.append(a.text)
                    if p['class'][0] == 'comments' and p['class'][1] == 'is-size-7':
                        comment = p.text.replace('\n', ' ').replace('\t', ' ').replace('\r', ' ')
                        
                result.append({
                    'arxiv': arxiv_info,
                    'title': title,
                    'author_info': author_info,
                    'summary': summary,
                    'comment': comment
                })
                if len(result) % 20 == 0:
                    print('Survey Papers Number: ', len(result))
                    
        start = len(result)
        if len(result) >= max_papers or len(result) == sum_num:
            break
        
    driver.quit()
    
    with open(f'.\\save\\json\\{query}_arxiv.json', 'w') as file:
        js.dump({'data': result}, file)

    make_pptx(query, result)

def main():
    query_list = [
        'prompt',
        'visual question answering',
        'visual generation',
        'bert',
        'transformer',
        'multimodal',
        'procedural',
        'pre-train',
        'pretrain',
        'visual language'
    ]
    search_type = 'title'
    for query in query_list:
        arxiv_survey(query, search_type=search_type)


if __name__ == '__main__':
    main()
